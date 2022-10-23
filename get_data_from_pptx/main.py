from pptx import Presentation
import pandas as pd

data = []

# load presentation to get data from.
prs = Presentation('demo_slide.pptx')


# get data from table on the first slide.
slide = prs.slides[0] 
table = slide.shapes[2].table 
for line in table.rows:
        each = ""
        for bits in line.cells:
            each += bits.text_frame.text + ","
        data.append(each)


# save data from PPTX to file
with open('data_out.csv', 'w') as f:
    for each in data:
        f.write(each)
        f.write("\n")


# Load CSV file into Pandas DF(+ tidy up mess).    
df = pd.read_csv('data_out.csv')
df = df.dropna(axis=1, how='all')



# Display contents of DataFrame.
print(df.to_string(index=None))
